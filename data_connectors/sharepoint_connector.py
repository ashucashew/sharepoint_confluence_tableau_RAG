import asyncio
from typing import List, Dict, Any, Optional
from datetime import datetime
import aiohttp
import os
from pathlib import Path
import mimetypes
import io
import tempfile

# Office file parsing libraries
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
import PyPDF2

from .base_connector import BaseConnector, Document
from config import settings
from loguru import logger

class SharePointConnector(BaseConnector):
    """Connector for Microsoft SharePoint"""
    
    def __init__(self, config: Dict[str, Any]):
        super().__init__(config)
        self.session = None
        self.base_url = settings.sharepoint_url
        self.site_url = settings.sharepoint_site_url
        
    async def connect(self) -> bool:
        """Establish connection to SharePoint"""
        try:
            self.session = aiohttp.ClientSession(
                auth=aiohttp.BasicAuth(
                    settings.sharepoint_username,
                    settings.sharepoint_password
                ),
                headers={
                    'Accept': 'application/json',
                    'Content-Type': 'application/json'
                }
            )
            
            # Test connection by trying to access the site
            async with self.session.get(f"{self.base_url}/_api/web") as response:
                if response.status == 200:
                    logger.info(f"Successfully connected to SharePoint at {self.base_url}")
                    return True
                else:
                    logger.error(f"Failed to connect to SharePoint: {response.status}")
                    return False
                    
        except Exception as e:
            logger.error(f"Failed to connect to SharePoint: {e}")
            return False
    
    async def fetch_documents(self) -> List[Document]:
        """Fetch all documents from SharePoint"""
        documents = []
        
        try:
            # Get all document libraries
            libraries = await self._get_document_libraries()
            
            for library in libraries:
                try:
                    library_docs = await self._fetch_library_documents(library)
                    documents.extend(library_docs)
                except Exception as e:
                    logger.error(f"Error fetching documents from library {library.get('name')}: {e}")
                    continue
            
            logger.info(f"Fetched {len(documents)} documents from SharePoint")
            return documents
            
        except Exception as e:
            logger.error(f"Error fetching SharePoint documents: {e}")
            return []
    
    async def _get_document_libraries(self) -> List[Dict[str, Any]]:
        """Get all document libraries from the site"""
        try:
            url = f"{self.site_url}/_api/web/lists?$filter=BaseTemplate eq 101"
            async with self.session.get(url) as response:
                if response.status == 200:
                    data = await response.json()
                    return data.get('value', [])
                else:
                    logger.error(f"Failed to get document libraries: {response.status}")
                    return []
        except Exception as e:
            logger.error(f"Error getting document libraries: {e}")
            return []
    
    async def _fetch_library_documents(self, library: Dict[str, Any]) -> List[Document]:
        """Fetch documents from a specific library"""
        documents = []
        
        try:
            library_name = library.get('Title', library.get('name', 'Unknown'))
            library_id = library.get('Id', library.get('id'))
            
            # Get all items in the library
            url = f"{self.site_url}/_api/web/lists(guid'{library_id}')/items?$expand=File,Author,Editor&$top=1000"
            
            async with self.session.get(url) as response:
                if response.status == 200:
                    data = await response.json()
                    items = data.get('value', [])
                    
                    for item in items:
                        try:
                            doc = await self._process_library_item(item, library_name)
                            if doc:
                                documents.append(doc)
                        except Exception as e:
                            logger.error(f"Error processing item {item.get('Id')}: {e}")
                            continue
                else:
                    logger.error(f"Failed to get items from library {library_name}: {response.status}")
                    
        except Exception as e:
            logger.error(f"Error fetching documents from library {library.get('name')}: {e}")
        
        return documents
    
    async def _process_library_item(self, item: Dict[str, Any], library_name: str) -> Optional[Document]:
        """Process a library item into a Document"""
        try:
            file_info = item.get('File', {})
            if not file_info:
                return None
            
            file_name = file_info.get('Name', '')
            file_url = file_info.get('ServerRelativeUrl', '')
            file_size = file_info.get('Length', 0)
            
            # Check if file type is supported
            file_extension = Path(file_name).suffix.lower()
            if file_extension not in settings.supported_file_types:
                return None
            
            # Get file content based on type
            content = await self._extract_file_content(file_url, file_extension)
            
            # Get metadata
            metadata = {
                'library_name': library_name,
                'file_size': file_size,
                'file_type': file_extension,
                'author': item.get('Author', {}).get('Title', 'Unknown'),
                'editor': item.get('Editor', {}).get('Title', 'Unknown'),
                'created': item.get('Created', ''),
                'modified': item.get('Modified', ''),
                'file_url': file_url
            }
            
            # Parse dates
            created_at = None
            last_modified = None
            
            try:
                if item.get('Created'):
                    created_at = datetime.fromisoformat(item['Created'].replace('Z', '+00:00'))
                if item.get('Modified'):
                    last_modified = datetime.fromisoformat(item['Modified'].replace('Z', '+00:00'))
            except:
                pass
            
            return Document(
                id=f"sharepoint_{item.get('Id')}",
                title=file_name,
                content=content,
                source_type="sharepoint",
                source_url=f"{self.base_url}{file_url}",
                metadata=metadata,
                last_modified=last_modified,
                created_at=created_at
            )
            
        except Exception as e:
            logger.error(f"Error processing library item {item.get('Id')}: {e}")
            return None
    
    async def _extract_file_content(self, file_url: str, file_extension: str) -> str:
        """Extract content from a file based on its type"""
        try:
            if file_extension in ['.txt', '.md', '.html']:
                # Download and read text files
                return await self._extract_text_content(file_url)
            
            elif file_extension == '.docx':
                # Extract content from Word documents
                return await self._extract_word_content(file_url)
            
            elif file_extension == '.pptx':
                # Extract content from PowerPoint presentations
                return await self._extract_powerpoint_content(file_url)
            
            elif file_extension == '.xlsx':
                # Extract content from Excel spreadsheets
                return await self._extract_excel_content(file_url)
            
            elif file_extension == '.pdf':
                # Extract content from PDF files
                return await self._extract_pdf_content(file_url)
            
            else:
                return f"File: {Path(file_url).name}"
                
        except Exception as e:
            logger.error(f"Error extracting content from {file_url}: {e}")
            return f"Error reading file: {Path(file_url).name}"
    
    async def _download_file(self, file_url: str) -> bytes:
        """Download file content from SharePoint"""
        full_url = f"{self.base_url}{file_url}"
        async with self.session.get(full_url) as response:
            if response.status == 200:
                return await response.read()
            else:
                raise Exception(f"Failed to download file: {response.status}")
    
    async def _extract_text_content(self, file_url: str) -> str:
        """Extract content from text files"""
        try:
            full_url = f"{self.base_url}{file_url}"
            async with self.session.get(full_url) as response:
                if response.status == 200:
                    return await response.text()
                else:
                    return f"Error downloading text file: {response.status}"
        except Exception as e:
            logger.error(f"Error extracting text content: {e}")
            return f"Error reading text file: {Path(file_url).name}"
    
    async def _extract_word_content(self, file_url: str) -> str:
        """Extract text content from Word documents"""
        try:
            # Download the file
            file_content = await self._download_file(file_url)
            
            # Parse Word document
            doc = DocxDocument(io.BytesIO(file_content))
            
            # Extract text from paragraphs
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
            
            # Extract text from tables
            tables = []
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                if table_text:
                    tables.append("TABLE:\n" + "\n".join(table_text) + "\nEND TABLE")
            
            # Combine paragraphs and tables
            content_parts = []
            if paragraphs:
                content_parts.append("\n\n".join(paragraphs))
            if tables:
                content_parts.extend(tables)
            
            return "\n\n".join(content_parts) if content_parts else "Empty Word document"
            
        except Exception as e:
            logger.error(f"Error extracting Word content: {e}")
            return f"Error reading Word document: {Path(file_url).name}"
    
    async def _extract_powerpoint_content(self, file_url: str) -> str:
        """Extract text content from PowerPoint presentations"""
        try:
            # Download the file
            file_content = await self._download_file(file_url)
            
            # Parse PowerPoint presentation
            prs = Presentation(io.BytesIO(file_content))
            
            # Extract content from each slide
            slides_content = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"Slide {slide_num}:"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                # Extract text from text boxes
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text.strip():
                                slide_text.append(paragraph.text.strip())
                
                # Extract text from tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_text.append("TABLE:\n" + "\n".join(table_text) + "\nEND TABLE")
                
                if len(slide_text) > 1:  # More than just the slide number
                    slides_content.append("\n".join(slide_text))
            
            return "\n\n".join(slides_content) if slides_content else "Empty PowerPoint presentation"
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint content: {e}")
            return f"Error reading PowerPoint presentation: {Path(file_url).name}"
    
    async def _extract_excel_content(self, file_url: str) -> str:
        """Extract content from Excel spreadsheets"""
        try:
            # Download the file
            file_content = await self._download_file(file_url)
            
            # Load Excel workbook
            wb = load_workbook(io.BytesIO(file_content), data_only=True)
            
            # Extract content from each worksheet
            worksheets_content = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_text = [f"Worksheet: {sheet_name}"]
                
                # Get the used range
                max_row = ws.max_row
                max_col = ws.max_column
                
                if max_row > 0 and max_col > 0:
                    # Extract data from cells
                    for row in range(1, min(max_row + 1, 100)):  # Limit to first 100 rows
                        row_data = []
                        for col in range(1, min(max_col + 1, 20)):  # Limit to first 20 columns
                            cell_value = ws.cell(row=row, column=col).value
                            if cell_value is not None:
                                row_data.append(str(cell_value).strip())
                        if row_data:
                            sheet_text.append(" | ".join(row_data))
                
                if len(sheet_text) > 1:  # More than just the worksheet name
                    worksheets_content.append("\n".join(sheet_text))
            
            return "\n\n".join(worksheets_content) if worksheets_content else "Empty Excel workbook"
            
        except Exception as e:
            logger.error(f"Error extracting Excel content: {e}")
            return f"Error reading Excel workbook: {Path(file_url).name}"
    
    async def _extract_pdf_content(self, file_url: str) -> str:
        """Extract text content from PDF files"""
        try:
            # Download the file
            file_content = await self._download_file(file_url)
            
            # Parse PDF
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            
            # Extract text from each page
            pages_content = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    pages_content.append(f"Page {page_num}:\n{text.strip()}")
            
            return "\n\n".join(pages_content) if pages_content else "Empty PDF document"
            
        except Exception as e:
            logger.error(f"Error extracting PDF content: {e}")
            return f"Error reading PDF document: {Path(file_url).name}"
    
    async def fetch_recent_documents(self, since: datetime) -> List[Document]:
        """Fetch documents modified since a specific time"""
        all_docs = await self.fetch_documents()
        recent_docs = [
            doc for doc in all_docs 
            if doc.last_modified and doc.last_modified > since
        ]
        logger.info(f"Found {len(recent_docs)} recent documents since {since}")
        return recent_docs
    
    async def get_document_by_id(self, doc_id: str) -> Optional[Document]:
        """Fetch a specific document by ID"""
        try:
            if doc_id.startswith("sharepoint_"):
                item_id = doc_id.replace("sharepoint_", "")
                # This would require implementing a specific item fetch
                # For now, we'll return None
                logger.warning("Individual document fetch not implemented for SharePoint")
                return None
        except Exception as e:
            logger.error(f"Error fetching document {doc_id}: {e}")
        return None
    
    async def __aenter__(self):
        return self
    
    async def __aexit__(self, exc_type, exc_val, exc_tb):
        if self.session:
            await self.session.close()
    
    def __del__(self):
        if self.session and not self.session.closed:
            asyncio.create_task(self.session.close()) 